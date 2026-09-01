import io
import json
import csv
import logging
from typing import List, Dict, Any

logger = logging.getLogger(__name__)

# Trilingual OCR per the product's English/Hindi/Marathi coverage — requires the
# tesseract-ocr-hin and tesseract-ocr-mar language packs (see backend/Dockerfile).
TESSERACT_LANG = "eng+hin+mar"

# T90 — lazily-loaded singleton so the model loads once per process, not
# once per page. lang='mr' selects PaddleOCR's Devanagari-family
# recognition model (covers Marathi and Hindi both — see paddleocr's
# DEVANAGARI_LANGS grouping), matching this product's stated priority
# ("the product cannot currently read the script it is sold on").
# Tradeoff versus Tesseract's eng+hin+mar single pass: PaddleOCR needs one
# language pipeline per call, so pure-English pages may see lower
# accuracy here than with Tesseract's combined pass — not a bug, a
# genuine engine tradeoff, left for a caller to choose via ocr_engine.
_paddle_ocr_instance = None


def _get_paddle_ocr():
    global _paddle_ocr_instance
    if _paddle_ocr_instance is None:
        from paddleocr import PaddleOCR
        _paddle_ocr_instance = PaddleOCR(
            lang="mr",
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            # Required on this platform — PaddlePaddle 3.3.1's default PIR
            # executor errors on the oneDNN-fused instruction path for the
            # text detection model:
            # "NotImplementedError: ConvertPirAttribute2RuntimeAttribute
            #  not support [pir::ArrayAttribute<pir::DoubleAttribute>]"
            # Confirmed live: identical inputs succeed with this off,
            # fail with it on. Revisit if a future paddlepaddle release
            # fixes the oneDNN/PIR interaction.
            enable_mkldnn=False,
        )
    return _paddle_ocr_instance


def _paddle_ocr_image(pil_img) -> str:
    import numpy as np
    ocr = _get_paddle_ocr()
    arr = np.array(pil_img.convert("RGB"))
    lines = []
    for res in ocr.predict(arr):
        lines.extend(res.get("rec_texts", []))
    return "\n".join(lines)


def extract_pages_from_file(file_bytes: bytes, filename: str, ocr_engine: str = "tesseract") -> List[Dict[str, Any]]:
    """
    Extract structured pages and text content from various file formats:
    PDF, Word (.docx), Excel (.xlsx, .csv), PowerPoint (.pptx), Markdown (.md),
    RTF (.rtf), JSON (.json), Images (.jpg, .png, etc.), and Plain Text files.

    ocr_engine ('tesseract' | 'paddle') only affects the PDF/image paths —
    every other format here doesn't involve OCR at all.
    """
    ext = filename.lower().split(".")[-1] if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(file_bytes, filename, ocr_engine)
    elif ext in ["docx", "doc"]:
        return _extract_docx(file_bytes)
    elif ext in ["xlsx", "xls"]:
        return _extract_excel(file_bytes)
    elif ext in ["pptx", "ppt"]:
        return _extract_pptx(file_bytes)
    elif ext == "csv":
        return _extract_csv(file_bytes)
    elif ext == "rtf":
        return _extract_rtf(file_bytes)
    elif ext == "json":
        return _extract_json(file_bytes)
    elif ext in ["jpg", "jpeg", "png", "bmp", "webp", "tiff"]:
        return _extract_image(file_bytes, filename, ocr_engine)
    else:
        return _extract_text(file_bytes)


def _extract_image(file_bytes: bytes, filename: str, ocr_engine: str = "tesseract") -> List[Dict[str, Any]]:
    text = ""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        if ocr_engine == "paddle":
            text = _paddle_ocr_image(img) or ""
        else:
            import pytesseract
            text = pytesseract.image_to_string(img, lang=TESSERACT_LANG) or ""
        logger.info(f"Image OCR ({ocr_engine}) extracted {len(text)} chars from {filename}")
    except Exception as e:
        logger.warning(f"Failed to perform {ocr_engine} OCR on image file {filename}: {e}")

    failed = not text.strip()
    if failed:
        text = f"Image document: {filename}"

    return [{
        "page_number": 1,
        "text": text.strip(),
        "words": [],
        "bbox": {},
        "extraction_failed": failed
    }]


def _extract_pdf(file_bytes: bytes, filename: str = "", ocr_engine: str = "tesseract") -> List[Dict[str, Any]]:
    import pdfplumber
    pages = []
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                raw_words = page.extract_words() or []

                # T05 — carry word-level regions from the extractor to the
                # fact writer instead of discarding them here. Normalised
                # to 0-1 page fractions, top-left origin, per T06's signed
                # coordinate contract (pdfplumber's top/bottom are already
                # top-down, so no y-flip needed). Page-level granularity
                # (every word on the page, not sliced per-chunk) — chunking
                # re-tokenizes page text and doesn't preserve a reliable
                # char-offset mapping back to pdfplumber's word list, so a
                # per-chunk slice would be a guess; the full per-page list
                # is exact and still a real, usable region source (T04's
                # own FactRegion works at page granularity too).
                pw = float(page.width) or 1.0
                ph = float(page.height) or 1.0
                words = [
                    {
                        "text": w.get("text", ""),
                        "x0": max(0.0, min(1.0, w["x0"] / pw)),
                        "y0": max(0.0, min(1.0, w["top"] / ph)),
                        "x1": max(0.0, min(1.0, w["x1"] / pw)),
                        "y1": max(0.0, min(1.0, w["bottom"] / ph)),
                    }
                    for w in raw_words
                    if "x0" in w and "top" in w and "x1" in w and "bottom" in w
                ]

                # OCR Fallback for scanned/image PDF pages
                if not text.strip():
                    try:
                        pil_img = page.to_image(resolution=150).original
                        if ocr_engine == "paddle":
                            ocr_text = _paddle_ocr_image(pil_img) or ""
                        else:
                            import pytesseract
                            ocr_text = pytesseract.image_to_string(pil_img, lang=TESSERACT_LANG) or ""
                        if ocr_text.strip():
                            text = ocr_text
                            logger.info(f"{ocr_engine} OCR extracted {len(text)} chars from page {i+1} of {filename}")
                    except Exception as ocr_err:
                        logger.warning(f"OCR fallback ({ocr_engine}) failed for page {i+1} of {filename}: {ocr_err}")

                failed = False
                if not text.strip():
                    text = f"Scanned page {i+1} of document {filename}"
                    failed = True

                pages.append({
                    "page_number": i + 1,
                    "text": text.strip(),
                    "words": words,
                    "bbox": {"width": float(page.width), "height": float(page.height)},
                    "extraction_failed": failed
                })
    except Exception as e:
        logger.error(f"Error parsing PDF with pdfplumber: {e}")
        pages.append({
            "page_number": 1,
            "text": f"Scanned PDF document: {filename}",
            "words": [],
            "bbox": {},
            "extraction_failed": True
        })
    return pages if pages else [{"page_number": 1, "text": f"Scanned document: {filename}", "words": [], "bbox": {}, "extraction_failed": True}]


def _clean_binary_strings(file_bytes: bytes) -> str:
    import re
    raw_str = file_bytes.decode("latin-1", errors="ignore")
    strings = re.findall(r'[\x20-\x7E\t\n\r]{3,}', raw_str)
    cleaned = [s.strip() for s in strings if len(s.strip()) >= 3 and not s.strip().startswith(('PK', 'Root Entry', '\x00'))]
    return "\n".join(cleaned)


def _extract_docx(file_bytes: bytes) -> List[Dict[str, Any]]:
    pages = []
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        full_text_blocks = []
        
        # Extract paragraph text
        for p in doc.paragraphs:
            if p.text.strip():
                full_text_blocks.append(p.text.strip())

        # Extract tables text
        for table in doc.tables:
            table_lines = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells]
                table_lines.append(" | ".join(row_cells))
            if table_lines:
                full_text_blocks.append("\n".join(table_lines))

        full_text = "\n\n".join(full_text_blocks)
        
        # Chunk text into ~1500 character logical pages for downstream chunker
        chunk_size = 1500
        paragraphs = full_text.split("\n\n")
        current_page_text = []
        current_len = 0
        page_num = 1

        for p in paragraphs:
            current_page_text.append(p)
            current_len += len(p)
            if current_len >= chunk_size:
                pages.append({
                    "page_number": page_num,
                    "text": "\n\n".join(current_page_text),
                    "words": [],
                    "bbox": {}
                })
                page_num += 1
                current_page_text = []
                current_len = 0

        if current_page_text:
            pages.append({
                "page_number": page_num,
                "text": "\n\n".join(current_page_text),
                "words": [],
                "bbox": {}
            })
    except Exception as e:
        logger.error(f"Error parsing DOCX/DOC: {e}")
        extracted_text = _clean_binary_strings(file_bytes)
        pages.append({
            "page_number": 1,
            "text": extracted_text if extracted_text.strip() else "Document content extracted.",
            "words": [],
            "bbox": {}
        })

    return pages if pages else [{"page_number": 1, "text": "Document content", "words": [], "bbox": {}}]


def _extract_excel(file_bytes: bytes) -> List[Dict[str, Any]]:
    pages = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        
        for idx, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            sheet_lines = [f"Sheet: {sheet_name}"]
            
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(val) if val is not None else "" for val in row]
                if any(row_vals):
                    sheet_lines.append(" | ".join(row_vals))

            sheet_text = "\n".join(sheet_lines)
            if sheet_text.strip():
                pages.append({
                    "page_number": idx + 1,
                    "text": sheet_text,
                    "words": [],
                    "bbox": {}
                })
    except Exception as e:
        logger.error(f"Error parsing Excel file: {e}")
        # Try pandas fallback
        try:
            import pandas as pd
            excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
            for idx, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text = f"Sheet: {sheet_name}\n" + df.to_string()
                pages.append({
                    "page_number": idx + 1,
                    "text": text,
                    "words": [],
                    "bbox": {}
                })
        except Exception as inner_e:
            logger.error(f"Pandas Excel fallback failed: {inner_e}")
            pages.append({
                "page_number": 1,
                "text": file_bytes.decode("utf-8", errors="ignore"),
                "words": [],
                "bbox": {}
            })

    return pages if pages else [{"page_number": 1, "text": "Excel sheet content", "words": [], "bbox": {}}]


def _extract_pptx(file_bytes: bytes) -> List[Dict[str, Any]]:
    pages = []
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        
        for idx, slide in enumerate(prs.slides):
            slide_text_blocks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_blocks.append(shape.text.strip())
            
            slide_text = f"Slide {idx + 1}:\n" + "\n".join(slide_text_blocks)
            pages.append({
                "page_number": idx + 1,
                "text": slide_text,
                "words": [],
                "bbox": {}
            })
    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        pages.append({
            "page_number": 1,
            "text": file_bytes.decode("utf-8", errors="ignore"),
            "words": [],
            "bbox": {}
        })

    return pages if pages else [{"page_number": 1, "text": "Presentation content", "words": [], "bbox": {}}]


def _extract_csv(file_bytes: bytes) -> List[Dict[str, Any]]:
    pages = []
    try:
        content = file_bytes.decode("utf-8", errors="ignore")
        reader = csv.reader(content.splitlines())
        lines = []
        for row in reader:
            if any(row):
                lines.append(" | ".join(row))

        # Split into pages of 100 rows each
        page_size = 100
        for i in range(0, max(1, len(lines)), page_size):
            chunk = lines[i:i+page_size]
            pages.append({
                "page_number": (i // page_size) + 1,
                "text": "\n".join(chunk),
                "words": [],
                "bbox": {}
            })
    except Exception as e:
        logger.error(f"Error parsing CSV: {e}")
        pages.append({
            "page_number": 1,
            "text": file_bytes.decode("utf-8", errors="ignore"),
            "words": [],
            "bbox": {}
        })

    return pages if pages else [{"page_number": 1, "text": "CSV document content", "words": [], "bbox": {}}]


def _extract_rtf(file_bytes: bytes) -> List[Dict[str, Any]]:
    try:
        from striprtf.striprtf import rtf_to_text
        raw_rtf = file_bytes.decode("utf-8", errors="ignore")
        plain_text = rtf_to_text(raw_rtf)
        return [{
            "page_number": 1,
            "text": plain_text.strip(),
            "words": [],
            "bbox": {}
        }]
    except Exception as e:
        logger.error(f"Error parsing RTF: {e}")
        return [{
            "page_number": 1,
            "text": file_bytes.decode("utf-8", errors="ignore"),
            "words": [],
            "bbox": {}
        }]


def _extract_json(file_bytes: bytes) -> List[Dict[str, Any]]:
    try:
        raw_text = file_bytes.decode("utf-8", errors="ignore")
        parsed_json = json.loads(raw_text)
        formatted_json = json.dumps(parsed_json, indent=2)
        return [{
            "page_number": 1,
            "text": formatted_json,
            "words": [],
            "bbox": {}
        }]
    except Exception as e:
        logger.error(f"Error parsing JSON: {e}")
        return [{
            "page_number": 1,
            "text": file_bytes.decode("utf-8", errors="ignore"),
            "words": [],
            "bbox": {}
        }]


def _extract_text(file_bytes: bytes) -> List[Dict[str, Any]]:
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode("latin-1")
        except Exception:
            text = file_bytes.decode("utf-8", errors="ignore")

    return [{
        "page_number": 1,
        "text": text if text.strip() else "Document text",
        "words": [],
        "bbox": {}
    }]
